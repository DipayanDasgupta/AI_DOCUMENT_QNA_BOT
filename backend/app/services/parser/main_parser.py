import os
import uuid
from typing import List, Dict, Any, Tuple, Optional
import io # For pandas reading in-memory
import json

# Parsing Libraries
from pypdf import PdfReader
import docx
from pptx import Presentation # For PPTX
import pandas as pd # For CSV, XLSX, JSON

# OCR & Image Handling
from PIL import Image
from pdf2image import convert_from_path
from pdf2image.exceptions import PDFInfoNotInstalledError

# Local Imports
from app.models.data_models import DocumentChunk
from app.services.parser.ocr import perform_ocr
from app.services.text_splitter import chunk_text, extract_urls
from app.core.config import settings

# --- PDF Parsing ---
async def _parse_pdf(file_path: str, original_name: str, session_id: str) -> Tuple[List[DocumentChunk], List[str]]:
    chunks = []
    urls = set()
    poppler_available = True # Assume available unless check fails
    try:
        # Optional: Poppler check (can be slow)
        # try: convert_from_path(file_path, first_page=1, last_page=1); print("[Parser] Poppler available.")
        # except Exception: poppler_available = False; print("[Parser] WARNING: Poppler not found/working. OCR fallback skipped.")
        pass # Skip check for speed for now

        reader = PdfReader(file_path)
        num_pages = len(reader.pages)
        print(f"[Parser] Processing PDF: {original_name}, Pages: {num_pages}")
        for page_num in range(num_pages):
            page_content = ""; ocr_attempted = False
            try:
                page = reader.pages[page_num]
                page_text = page.extract_text()
                if page_text and page_text.strip(): page_content = page_text
                elif poppler_available: # Try OCR only if text extraction fails and poppler *might* be there
                    print(f"[Parser] Page {page_num+1}: No text found, attempting OCR...")
                    ocr_attempted = True
                    images = convert_from_path(file_path, dpi=200, first_page=page_num + 1, last_page=page_num + 1, fmt='jpeg', thread_count=1)
                    if images: ocr_text = perform_ocr(images[0]); page_content = ocr_text or ""
                    if page_content: print(f"[Parser] Page {page_num+1}: OCR successful.")
                    else: print(f"[Parser] Page {page_num+1}: OCR yielded no text.")
                else: print(f"[Parser] Page {page_num+1}: No text found, OCR skipped (Poppler unavailable?).")

            except PDFInfoNotInstalledError: # Catch specific poppler error here
                 print("[Parser] WARNING: Poppler not installed - Skipping OCR attempt for PDF page.")
                 poppler_available = False # Don't try again for this PDF
            except Exception as page_err: print(f"[Parser] Error processing page {page_num+1}: {page_err}")

            if page_content:
                page_urls = extract_urls(page_content); urls.update(page_urls)
                split_chunks = chunk_text(page_content)
                for text_chunk in split_chunks:
                    chunks.append(DocumentChunk(session_id=session_id, source=original_name, text=text_chunk, page=page_num + 1, metadata={"parser": "pdf", "ocr_attempted": ocr_attempted}))
    except Exception as e: print(f"[Parser] ERROR Failed to process PDF {original_name}: {e}"); traceback.print_exc()
    return chunks, list(urls)

# --- DOCX Parsing ---
async def _parse_docx(file_path: str, original_name: str, session_id: str) -> Tuple[List[DocumentChunk], List[str]]:
    chunks = []; urls = set()
    try:
        doc = docx.Document(file_path); print(f"[Parser] Processing DOCX: {original_name}")
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        # TODO: Could potentially extract text from tables as well
        if full_text:
            urls.update(extract_urls(full_text))
            split_chunks = chunk_text(full_text)
            for text_chunk in split_chunks: chunks.append(DocumentChunk(session_id=session_id, source=original_name, text=text_chunk, metadata={"parser": "docx"}))
    except Exception as e: print(f"[Parser] ERROR Failed to read DOCX {original_name}: {e}")
    return chunks, list(urls)

# --- PPTX Parsing ---
async def _parse_pptx(file_path: str, original_name: str, session_id: str) -> Tuple[List[DocumentChunk], List[str]]:
    chunks = []; urls = set()
    try:
        prs = Presentation(file_path); print(f"[Parser] Processing PPTX: {original_name}")
        full_text = ""
        for i, slide in enumerate(prs.slides):
            slide_text = f"--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            # Also check notes slide
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame and slide.notes_slide.notes_text_frame.text.strip():
                 slide_text += "\nNotes:\n" + slide.notes_slide.notes_text_frame.text + "\n"
            full_text += slide_text + "\n"

        if full_text:
            urls.update(extract_urls(full_text))
            split_chunks = chunk_text(full_text) # Chunk the entire presentation text
            for text_chunk in split_chunks: chunks.append(DocumentChunk(session_id=session_id, source=original_name, text=text_chunk, metadata={"parser": "pptx"}))
    except Exception as e: print(f"[Parser] ERROR Failed to read PPTX {original_name}: {e}")
    return chunks, list(urls)

# --- TXT Parsing ---
async def _parse_txt(file_path: str, original_name: str, session_id: str) -> Tuple[List[DocumentChunk], List[str]]:
    chunks = []; urls = set()
    try:
        print(f"[Parser] Processing TXT: {original_name}")
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f: full_text = f.read()
        if full_text:
            urls.update(extract_urls(full_text))
            split_chunks = chunk_text(full_text)
            for text_chunk in split_chunks: chunks.append(DocumentChunk(session_id=session_id, source=original_name, text=text_chunk, metadata={"parser": "text"}))
    except Exception as e: print(f"[Parser] ERROR Failed to read TXT {original_name}: {e}")
    return chunks, list(urls)

# --- Image Parsing (OCR) ---
async def _parse_image(file_path: str, original_name: str, session_id: str) -> Tuple[List[DocumentChunk], List[str]]:
    chunks = []; urls = set()
    print(f"[Parser] Processing Image for OCR: {original_name}")
    ocr_text = perform_ocr(file_path)
    if ocr_text:
        urls.update(extract_urls(ocr_text))
        split_chunks = chunk_text(ocr_text)
        for text_chunk in split_chunks: chunks.append(DocumentChunk(session_id=session_id, source=original_name, text=text_chunk, metadata={"parser": "ocr"}))
    else: print(f"[Parser] No text extracted via OCR from {original_name}")
    return chunks, list(urls)

# --- Structured Data Parsing (CSV, XLSX, JSON) ---
def generate_df_summary(df: pd.DataFrame, filename: str) -> str:
    """Creates a textual summary of a pandas DataFrame."""
    summary = f"Summary for {filename}:\n"
    buffer = io.StringIO()
    df.info(buf=buffer)
    summary += "Schema/Info:\n" + buffer.getvalue() + "\n"
    summary += "First 5 Rows:\n" + df.head().to_string() + "\n"
    # Could add more: df.describe().to_string() for numerical stats
    return summary

async def _parse_structured(file_path: str, original_name: str, session_id: str) -> Tuple[List[DocumentChunk], List[str], Optional[pd.DataFrame]]:
    """Parses CSV, XLSX, JSON into DataFrame and generates summary."""
    chunks = []; urls = set(); df = None
    file_extension = os.path.splitext(original_name)[1].lower()
    print(f"[Parser] Processing Structured Data ({file_extension}): {original_name}")
    try:
        if file_extension == '.csv':
            df = pd.read_csv(file_path)
        elif file_extension == '.xlsx':
            df = pd.read_excel(file_path)
        elif file_extension == '.json':
            # Try loading common JSON structures (list of records or single object)
            try: df = pd.read_json(file_path, orient='records', lines=True) # Try lines first
            except ValueError: df = pd.read_json(file_path, orient='records') # Try standard records array
            except Exception: # Fallback for complex/nested JSON
                 with open(file_path, 'r') as f: raw_json_text = f.read()
                 # Just chunk the raw JSON text if pandas fails
                 split_chunks = chunk_text(raw_json_text); df=None
                 for text_chunk in split_chunks: chunks.append(DocumentChunk(session_id=session_id, source=original_name, text=text_chunk, metadata={"parser": "json_raw"}))
                 print("[Parser] Parsed JSON as raw text.")
                 return chunks, list(urls), df # Return early if raw parse

        if df is not None:
            summary_text = generate_df_summary(df, original_name)
            urls.update(extract_urls(summary_text)) # Unlikely but check summary
            split_chunks = chunk_text(summary_text) # Chunk the summary
            for text_chunk in split_chunks: chunks.append(DocumentChunk(session_id=session_id, source=original_name, text=text_chunk, metadata={"parser": f"structured_{file_extension}", "is_summary": True}))
            print(f"[Parser] Generated {len(split_chunks)} summary chunks for {original_name}.")

    except Exception as e:
        print(f"[Parser] ERROR Failed to read structured data {original_name}: {e}")
        df = None # Ensure df is None on error

    return chunks, list(urls), df


# --- Main Processing Function ---
# Add DataFrame storage to return value
async def process_document(file_path: str, original_name: str, session_id: str) -> Tuple[List[DocumentChunk], List[str], Optional[pd.DataFrame]]:
    """
    Routes file processing based on extension.
    Returns chunks, unique URLs, and DataFrame (for structured types).
    """
    import traceback # Ensure traceback is imported here too
    print(f"[Parser Service] Routing: {original_name} (Session: {session_id})")
    file_extension = os.path.splitext(original_name)[1].lower()
    chunks = []; urls = []; dataframe = None

    try:
        if file_extension == '.pdf':
            chunks, urls = await _parse_pdf(file_path, original_name, session_id)
        elif file_extension == '.docx':
             chunks, urls = await _parse_docx(file_path, original_name, session_id)
        elif file_extension == '.pptx':
             chunks, urls = await _parse_pptx(file_path, original_name, session_id)
        elif file_extension == '.txt':
             chunks, urls = await _parse_txt(file_path, original_name, session_id)
        elif file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
            chunks, urls = await _parse_image(file_path, original_name, session_id)
        elif file_extension in ['.csv', '.xlsx', '.json']:
             chunks, urls, dataframe = await _parse_structured(file_path, original_name, session_id)
        else:
            print(f"[Parser Service] Warning: Unsupported file type '{file_extension}' for {original_name}")

        unique_urls = sorted(list(set(urls))) # Ensure unique URLs
        print(f"[Parser Service] Finished processing {original_name}. Found {len(chunks)} chunks, {len(unique_urls)} unique URLs. DataFrame generated: {dataframe is not None}")

    except Exception as e:
         print(f"[Parser Service] CRITICAL ERROR during processing of {original_name}: {e}")
         traceback.print_exc()

    # Return DataFrame only for structured types
    return chunks, unique_urls, dataframe if file_extension in ['.csv', '.xlsx', '.json'] else None

